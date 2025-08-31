#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pptx（PowerPoint）からテキストを抽出し、1ファイル=1行のJSONとしてJSONLに保存するスクリプト。

- 入力: --input <dir>   （既定: ./pptx）
- 出力: --output <file> （既定: ./pptx_texts.jsonl）
- オプション: --include-notes （ノート欄のテキストも含める）
- 依存: python-pptx  (pip install python-pptx)

各PPTXにつき、抽出した全テキストをまとめて "text" フィールドに格納します。
"""

import argparse
import json
from pathlib import Path
from typing import List, Dict, Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def gather_pptx_files(input_dir: Path) -> List[Path]:
    """入力ディレクトリ直下の .pptx を列挙（拡張子の大小文字は無視）"""
    return sorted([p for p in input_dir.iterdir()
                   if p.is_file() and p.suffix.lower() == ".pptx"])


def extract_texts_from_shape(shape) -> List[str]:
    """単一の図形からテキストを抽出（表・グループ・テキストフレームを再帰的に処理）"""
    texts: List[str] = []
    st = shape.shape_type

    # グループ図形
    if st == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            texts.extend(extract_texts_from_shape(s))

    # 表
    elif st == MSO_SHAPE_TYPE.TABLE:
        tbl = shape.table
        for row in tbl.rows:
            for cell in row.cells:
                t = (cell.text or "").strip()
                if t:
                    texts.append(t)

    # 通常のテキストフレーム
    elif getattr(shape, "has_text_frame", False):
        t = (shape.text or "").strip()
        if t:
            texts.append(t)

    # それ以外は無視（画像など）
    return texts


def extract_texts_from_slide(slide) -> List[str]:
    """スライド内の全図形からテキストを抽出して平坦化"""
    out: List[str] = []
    for shape in slide.shapes:
        out.extend(extract_texts_from_shape(shape))
    # 空行を除去
    return [t for t in (s.strip() for s in out) if t]


def extract_notes_text(slide) -> str:
    """ノート欄テキスト（存在すれば）を返す"""
    try:
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            return (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        pass
    return ""


def parse_pptx_to_record(pptx_path: Path, include_notes: bool = False) -> Dict[str, Any]:
    """1つのPPTXを1つのJSONオブジェクト（辞書）に変換"""
    prs = Presentation(pptx_path)
    slides_json: List[Dict[str, Any]] = []
    all_text_chunks: List[str] = []

    for idx, slide in enumerate(prs.slides):
        slide_texts = extract_texts_from_slide(slide)
        slide_obj: Dict[str, Any] = {"index": idx, "texts": slide_texts}

        if include_notes:
            notes = extract_notes_text(slide)
            if notes:
                slide_obj["notes"] = notes

        slides_json.append(slide_obj)

        # 文書全体の連結用
        all_text_chunks.extend(slide_texts)
        if include_notes and "notes" in slide_obj:
            all_text_chunks.append(slide_obj["notes"])

    record: Dict[str, Any] = {
        "file": pptx_path.name,
        "num_slides": len(prs.slides),
        "text": "\n".join(all_text_chunks),  # ← まとめて "text" に格納
        "slides": slides_json                # （必要なければ後段で無視/削除可）
    }
    return record


def main():
    ap = argparse.ArgumentParser(description="Extract PPTX texts to JSONL (1 pptx -> 1 line).")
    ap.add_argument("--input", type=Path, default=Path("pptx"), help="入力ディレクトリ（既定: ./pptx）")
    ap.add_argument("--output", type=Path, default=Path("pptx_texts.jsonl"), help="出力JSONLファイル")
    ap.add_argument("--include-notes", action="store_true", help="ノート欄のテキストも含める")
    args = ap.parse_args()

    args.input.mkdir(exist_ok=True, parents=True)

    files = gather_pptx_files(args.input)
    if not files:
        print(f"[WARN] PPTX が見つかりませんでした: {args.input.resolve()}")
        return

    count_ok = 0
    with args.output.open("w", encoding="utf-8") as wf:
        for pptx_path in files:
            try:
                rec = parse_pptx_to_record(pptx_path, include_notes=args.include_notes)
            except Exception as e:
                # 失敗しても他のファイルは続行
                rec = {"file": pptx_path.name, "error": str(e)}
            wf.write(json.dumps(rec, ensure_ascii=False) + "\n")
            if "error" not in rec:
                count_ok += 1

    print(f"=== 完了: {count_ok}/{len(files)} ファイルを書き出し → {args.output}")


if __name__ == "__main__":
    main()
